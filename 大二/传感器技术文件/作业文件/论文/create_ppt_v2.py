# -*- coding: utf-8 -*-
"""Generate LiDAR-focused PPT: 原理→元器件→测量电路→应用"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"D:\sensorhomework\论文"
FIG = os.path.join(BASE, "figure")
FORMULA_DIR = os.path.join(BASE, "formulas", "png")
CIRCUIT_DIR = os.path.join(BASE, "circuits")

prs = Presentation()
SW = prs.slide_width = Inches(13.333)
SH = prs.slide_height = Inches(7.5)

# ── Color Palette ──
DARK   = RGBColor(0x0F, 0x24, 0x48)
MED    = RGBColor(0x1A, 0x56, 0x9E)
LIGHT  = RGBColor(0x3B, 0x8E, 0xD8)
ORANGE = RGBColor(0xE8, 0x6A, 0x17)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xC0, 0x39, 0x2B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF2, 0xF5)
DGRAY  = RGBColor(0x2C, 0x3E, 0x50)
BLACK  = RGBColor(0x00, 0x00, 0x00)
ACCLR  = RGBColor(0x88, 0xAA, 0xCC)

def add_dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

def add_light_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def add_title_bar(slide, title_text, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.7); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.text = title_text
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = WHITE
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), SW, Inches(0.04))
    acc.fill.solid(); acc.fill.fore_color.rgb = ORANGE
    acc.line.fill.background()
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12), Inches(0.4))
        tf2 = box.text_frame; p2 = tf2.paragraphs[0]
        p2.text = subtitle; p2.font.size = Pt(16); p2.font.color.rgb = LIGHT

def add_circuit(slide, name, left, top, width=None, height=Inches(2.0)):
    png_path = os.path.join(CIRCUIT_DIR, f"{name}.png")
    if not os.path.exists(png_path):
        return None
    if width is None:
        width = Inches(5)
    return slide.shapes.add_picture(png_path, left, top, width, height)

def add_formula(slide, name, left, top, width=None, height=Inches(0.45)):
    png_path = os.path.join(FORMULA_DIR, f"{name}.png")
    if not os.path.exists(png_path):
        return None
    if width is None:
        width = Inches(4)
    return slide.shapes.add_picture(png_path, left, top, width, height)

def add_chart(slide, name, left, top, width, height=None):
    path = os.path.join(FIG, name)
    h = height or width * 0.6
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height=h)

def add_textbox(slide, left, top, width, height, items, default_size=Pt(16)):
    """Add multi-paragraph textbox. items = list of dicts: {text, size, color, bold, level}"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item.get("text", "")
        p.font.size = item.get("size", default_size)
        p.font.color.rgb = item.get("color", DGRAY)
        p.font.bold = item.get("bold", False)
        p.level = item.get("level", 0)
        p.space_after = Pt(item.get("space", 6))
    return box

def add_page_num(slide, n, total=19):
    box = slide.shapes.add_textbox(SW - Inches(1.3), SH - Inches(0.45), Inches(1.1), Inches(0.35))
    tf = box.text_frame; p = tf.paragraphs[0]
    p.text = f"{n}/{total}"
    p.font.size = Pt(10); p.font.color.rgb = LIGHT; p.alignment = PP_ALIGN.RIGHT

def add_highlight_box(slide, left, top, width, height, text, bg=RGBColor(0xFF, 0xF3, 0xE0), border=ORANGE, size=Pt(14)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = border; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text; p.font.size = size; p.font.color.rgb = border; p.font.bold = True

def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(s)
for y, h, c in [(0, 0.08, ORANGE), (0.08, 0.04, MED), (0.12, 0.02, LIGHT)]:
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), SW, Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

box = s.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(2.0))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "激光雷达技术深度解析"
p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "从物理原理、核心元器件、精密测量电路到自动驾驶应用"
p2.font.size = Pt(24); p2.font.color.rgb = ORANGE; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)

box2 = s.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.3), Inches(1.5))
tf2 = box2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "---- 兼论摄像机-毫米波雷达-超声波雷达-IMU的多模态融合"
p3.font.size = Pt(16); p3.font.color.rgb = ACCLR; p3.alignment = PP_ALIGN.CENTER
p4 = tf2.add_paragraph()
p4.text = "广东工业大学  传感器技术与应用  课程调研报告"
p4.font.size = Pt(15); p4.font.color.rgb = ACCLR; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(16)
p5 = tf2.add_paragraph()
p5.text = "指导老师：翟老师  -  2026年5月"
p5.font.size = Pt(13); p5.font.color.rgb = ACCLR; p5.alignment = PP_ALIGN.CENTER; p5.space_before = Pt(6)

rb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.3), SW, Inches(0.2))
rb.fill.solid(); rb.fill.fore_color.rgb = MED; rb.line.fill.background()

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "报告大纲 ---- 以激光雷达为主线的技术调研")

items = [
    {"text": "第一部分  激光雷达原理与物理基础", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "ToF飞行时间法 - FMCW相干探测 - 光电效应 - SPAD Geiger模式 - 探测统计", "level": 1, "size": Pt(15), "space": 10},
    {"text": "第二部分  激光雷达核心元器件", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "VCSEL/EEL光源 - SPAD/SiPM/APD探测器 - MEMS微振镜 - 905nm vs 1550nm技术路线", "level": 1, "size": Pt(15), "space": 10},
    {"text": "第三部分  激光雷达精密测量电路", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "TIA跨阻放大器 - TDC时间数字转换器 - SPAD主动/被动淬灭电路 - 完整信号链", "level": 1, "size": Pt(15), "space": 10},
    {"text": "第四部分  应用、融合与职业规划", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "自动驾驶感知 - 多传感器融合(BEV+Transformer) - 市场前景 - 岗位薪资", "level": 1, "size": Pt(15), "space": 10},
    {"text": "附录  其他四大传感器技术概览", "bold": True, "color": DGRAY, "size": Pt(19), "space": 12},
    {"text": "CMOS摄像机 - 毫米波雷达 - 超声波雷达 - MEMS IMU", "level": 1, "size": Pt(15)},
]
add_textbox(s, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.8), items)
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(1.5), Inches(0.05), Inches(5.5))
acc.fill.solid(); acc.fill.fore_color.rgb = LIGHT; acc.line.fill.background()
add_page_num(s, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: Overview -- Why LiDAR
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第一章  激光雷达概述 ---- 为什么LiDAR是L3+自动驾驶的\"眼睛\"？")

items_L = [
    {"text": "激光雷达 LiDAR = Light Detection And Ranging", "bold": True, "color": MED, "size": Pt(19), "space": 14},
    {"text": "主动发射激光脉冲 → 测量回波时间/频率/相位 → 三维点云", "level": 1, "size": Pt(16), "space": 10},
    {"text": "", "size": Pt(6)},
    {"text": "五大传感器中，LiDAR是唯一能直接输出高精度3D几何信息的传感器", "bold": True, "color": ORANGE, "size": Pt(19), "space": 14},
    {"text": "", "size": Pt(6)},
    {"text": "核心性能指标", "bold": True, "color": DARK, "size": Pt(16), "space": 10},
    {"text": "测距范围: 150m@10%反射率 → 300m@80%反射率 (L3级要求)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "角分辨率: 0.05°-0.2° (H) × 0.05°-0.5° (V)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "帧率: 10-30Hz (实时建图定位需求)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "探测概率: Pd > 99.7% @ 100m (ISO 22737标准)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "点云密度: 每秒数十万至数百万点 (禾赛AT128: 1.53M pts/s)", "level": 1, "size": Pt(15)},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.5), items_L)

# Right: market charts
add_chart(s, "market_size.png", Inches(6.9), Inches(1.5), Inches(6.0), Inches(2.8))
add_chart(s, "market_share.png", Inches(6.9), Inches(4.0), Inches(6.0), Inches(2.8))
add_chart(s, "cost_breakdown.png", Inches(6.9), Inches(6.0), Inches(6.0), height=Inches(1.2))

add_page_num(s, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: LiDAR Principle (1) -- ToF & FMCW
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第二章  激光雷达原理（一）---- ToF飞行时间法与FMCW相干探测")

items_L = [
    {"text": "ToF（飞行时间法）★ 量产主流，>90%市场份额", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "发射纳秒级激光脉冲(τp=1-10ns) → 目标反射 → 接收器检测回波", "level": 1, "size": Pt(15), "space": 8},
    {"text": "核心公式: R = ½-c-Δt", "bold": True, "level": 1, "size": Pt(16), "color": ORANGE, "space": 8},
    {"text": "1cm精度需要67ps计时分辨率----远小于单一时钟周期(1ns→1GHz)", "level": 2, "size": Pt(14), "space": 8},
    {"text": "直接ToF(dToF): 单光子计数+直方图峰值 → 抗干扰强", "level": 1, "size": Pt(15), "space": 8},
    {"text": "间接ToF(iToF): 调制光相位差 → 适合近距离高精度", "level": 1, "size": Pt(15), "space": 10},
    {"text": "", "size": Pt(6)},
    {"text": "FMCW（调频连续波）★ 下一代技术路线", "bold": True, "color": GREEN, "size": Pt(19), "space": 12},
    {"text": "线性调频激光 → 回波与本振光学混频 → 拍频检测", "level": 1, "size": Pt(15), "space": 8},
    {"text": "三大核心优势: ①直接测速(via多普勒) ②抗阳光+抗串扰 ③散粒噪声极限灵敏度", "level": 1, "size": Pt(15), "color": GREEN, "space": 8},
    {"text": "当前瓶颈: kHz窄线宽可调谐激光器成本高, 2026-27年工程样机", "level": 1, "size": Pt(14), "color": ORANGE},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5), items_L)

# Right side: formulas
add_formula(s, "LiDAR_ToF", Inches(7.3), Inches(1.5), width=Inches(5.2), height=Inches(0.55))
add_formula(s, "LidarEq", Inches(7.3), Inches(2.2), width=Inches(5.2), height=Inches(0.55))

items_R = [
    {"text": "ToF信号模型:", "bold": True, "color": DARK, "size": Pt(15), "space": 8},
    {"text": "P_r = P_t × (ρ×A_r)/(π×R²) × η_atm × η_opt", "level": 1, "size": Pt(15), "space": 6},
    {"text": "P_t发射功率, ρ目标反射率, A_r接收口径面积", "level": 2, "size": Pt(11), "space": 8},
    {"text": "激光雷达方程:", "bold": True, "color": DARK, "size": Pt(15), "space": 8},
    {"text": "SNR ∝ P_t×ρ×A_r×τ_p / (R²×NEP)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "传输损耗∝1/R², 比雷达(∝1/R⁴)有本质优势", "level": 2, "size": Pt(11), "space": 8},
    {"text": "FMCW关键公式:", "bold": True, "color": DARK, "size": Pt(15), "space": 8},
    {"text": "f_IF = 2B×R/(c×T_chirp)", "level": 1, "size": Pt(15), "color": ORANGE, "space": 6},
    {"text": "f_d = 2v_r/λ (同时获取速度!)", "level": 1, "size": Pt(15), "color": ORANGE, "space": 6},
    {"text": "ΔR = c/(2B) (距离分辨率由带宽决定)", "level": 1, "size": Pt(15)},
]
add_textbox(s, Inches(7.0), Inches(2.9), Inches(6.0), Inches(4.3), items_R)
add_page_num(s, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: LiDAR Principle (2) -- Photoelectric Physics
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第二章  激光雷达原理（二）---- 内光电效应与单光子探测物理")

items_L = [
    {"text": "内光电效应 ---- 光子→电子的转换基础", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "光子能量 E_ph = hc/λ > 半导体禁带宽度 E_g → 电子从价带跃迁到导带", "level": 1, "size": Pt(16), "space": 8},
    {"text": "硅(Si): E_g≈1.12eV → λ_max≤1100nm → 905nm可用", "level": 1, "size": Pt(16), "space": 8},
    {"text": "InGaAs: E_g≈0.75eV → λ_max≤1650nm → 1550nm可用", "level": 1, "size": Pt(16), "space": 8},
    {"text": "光生载流子在耗尽层强电场(>10⁵V/cm)下加速漂移 → 光电流", "level": 1, "size": Pt(15), "space": 12},
    {"text": "探测器增益机制对比", "bold": True, "color": MED, "size": Pt(19), "space": 10},
    {"text": "PIN光电二极管: 无内部增益, R=0.4-0.8A/W (量子效率η≈60-90%)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "APD雪崩光电二极管: 线性模式, 碰撞电离增益 M=50-200", "level": 1, "size": Pt(15), "space": 6},
    {"text": "SPAD单光子雪崩二极管: Geiger模式(偏压>击穿电压V_BR), 增益→∞(单光子触发)", "level": 1, "size": Pt(15), "color": ORANGE, "space": 6},
    {"text": "SiPM硅光电倍增管: N(10³-10⁴)个SPAD微单元并联 → 光子数分辨", "level": 1, "size": Pt(15), "color": GREEN},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5), items_L)

# Right: SPAD formulas
add_formula(s, "SPAD_prob", Inches(7.3), Inches(1.5), width=Inches(5.5), height=Inches(0.5))
add_formula(s, "SiPM", Inches(7.3), Inches(2.2), width=Inches(5.5), height=Inches(0.5))
add_formula(s, "DarkI", Inches(7.3), Inches(2.9), width=Inches(5.5), height=Inches(0.5))

items_R = [
    {"text": "SPAD探测概率(非齐次泊松过程):", "bold": True, "color": DARK, "size": Pt(15), "space": 8},
    {"text": "P_det = 1 - exp(-∫[η_PDE×Φ_ph(t) + DCR]dt)", "level": 1, "size": Pt(15), "space": 8},
    {"text": "η_PDE = η_QE × FF × P_av ≈ 5%-30% (光子探测效率)", "level": 2, "size": Pt(14), "space": 8},
    {"text": "SiPM输出电流:", "bold": True, "color": DARK, "size": Pt(15), "space": 8},
    {"text": "I_SiPM = (N_fired/N_micro) × G_SPAD × q", "level": 1, "size": Pt(15), "space": 8},
    {"text": "线性区: N_fired << N_micro → I ∝ Φ_ph (光强)", "level": 2, "size": Pt(14), "space": 8},
    {"text": "暗电流挑战:", "bold": True, "color": ORANGE, "size": Pt(15), "space": 8},
    {"text": "I_dark ∝ T^(3/2)×exp(-E_g/2kT)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "温度每↑8-10°C → 暗计数翻倍！ 85°C恶化百倍", "level": 1, "size": Pt(15), "color": RED, "space": 6},
    {"text": "对策: TEC半导体制冷 + 时间符合滤波 + 工艺降低缺陷密度", "level": 1, "size": Pt(14), "color": GREEN},
]
add_textbox(s, Inches(7.0), Inches(3.5), Inches(6.0), Inches(3.8), items_R)
add_page_num(s, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: LiDAR Components (1) -- Transmitter
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第三章  激光雷达元器件（一）---- 发射端：从激光器到光束扫描")

items_L = [
    {"text": "半导体激光光源", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "EEL边发射激光器(Edge Emitting Laser)", "bold": True, "level": 1, "size": Pt(16), "space": 6},
    {"text": "光从芯片侧面出射 - 高功率(>100W峰值) - 905nm GaAs基", "level": 2, "size": Pt(15), "space": 4},
    {"text": "代表: Osram SPL系列 - Lumentum - ams OSRAM", "level": 2, "size": Pt(14), "space": 8},
    {"text": "VCSEL垂直腔面发射激光器 ★ 新兴主流", "bold": True, "level": 1, "size": Pt(16), "space": 6},
    {"text": "光从芯片表面垂直出射 → 晶圆级测试+二维阵列", "level": 2, "size": Pt(15), "space": 4},
    {"text": "多结VCSEL: 10层有源区堆叠 → 功率密度3-5×提升", "level": 2, "size": Pt(15), "space": 4},
    {"text": "优势: 低成本-低发散角-高可靠性-可二维集成", "level": 2, "size": Pt(15), "color": GREEN, "space": 10},
    {"text": "905nm vs 1550nm 两大技术路线竞争", "bold": True, "color": MED, "size": Pt(19), "space": 10},
    {"text": "905nm(GaAs): 成本低-Si探测器兼容-产业链成熟-主导车载量产", "level": 1, "size": Pt(15), "color": GREEN, "space": 6},
    {"text": "1550nm(InP): 人眼安全高功率-大气穿透好-需InGaAs探测器-成本高", "level": 1, "size": Pt(15), "space": 6},
    {"text": "Luminar(1550nm路线) vs 禾赛/速腾聚创(905nm路线)", "level": 1, "size": Pt(14), "color": DGRAY},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.8), items_L)

items_R = [
    {"text": "光束扫描方案对比", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "① MEMS微振镜 ★ 当前量产主流", "bold": True, "level": 1, "size": Pt(16), "space": 6},
    {"text": "单轴/双轴谐振微镜 → 改变光束出射方向", "level": 2, "size": Pt(15), "space": 4},
    {"text": "代表: 禾赛AT128(128线, MEMS一维扫描+转镜)", "level": 2, "size": Pt(14), "space": 4},
    {"text": "优势: 成本<机械式-可靠性高-量产成熟", "level": 2, "size": Pt(14), "color": GREEN, "space": 8},
    {"text": "② OPA光学相控阵 ★ 理想纯固态", "bold": True, "level": 1, "size": Pt(16), "space": 6},
    {"text": "N个光栅天线相位调控 → 无机械运动部件", "level": 2, "size": Pt(15), "space": 4},
    {"text": "瓶颈: FOV受限±20°-旁瓣抑制难-硅光工艺", "level": 2, "size": Pt(14), "color": ORANGE, "space": 8},
    {"text": "③ Flash面阵闪光 ★ 近距离纯固态", "bold": True, "level": 1, "size": Pt(16), "space": 6},
    {"text": "单脉冲照亮整个FOV → SPAD面阵接收", "level": 2, "size": Pt(15), "space": 4},
    {"text": "瓶颈: 功率密度不足-探测距离<50m", "level": 2, "size": Pt(14), "color": ORANGE, "space": 8},
    {"text": "④ 机械旋转式 → 逐渐退出车载前装市场", "bold": True, "level": 1, "size": Pt(16), "color": DGRAY, "space": 6},
    {"text": "⑤ 转镜+多线束 → 速腾聚创M1, 棱镜旋转扫描", "level": 1, "size": Pt(15)},
]
add_textbox(s, Inches(7.2), Inches(1.5), Inches(5.8), Inches(5.0), items_R)
add_chart(s, "905_vs_1550.png", Inches(7.0), Inches(6.2), Inches(6.0), height=Inches(1.0))
add_page_num(s, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: LiDAR Components (2) -- Receiver Detector Chain
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第三章  激光雷达元器件（二）---- 接收端：单光子探测器链")

items_all = [
    {"text": "光电探测器四大层级 ---- 从μA到单光子", "bold": True, "color": MED, "size": Pt(18), "space": 14},
    {"text": "", "size": Pt(4)},
    {"text": "Level 1: PIN光电二极管 ---- 无增益，基础光电转换", "bold": True, "color": DARK, "size": Pt(16), "space": 6},
    {"text": "R=0.4-0.8A/W, 带宽>1GHz, 结构简单，适用于近距离低精度场景。Si-PIN/InGaAs-PIN", "level": 1, "size": Pt(15), "space": 8},
    {"text": "Level 2: APD雪崩光电二极管 ---- 线性模式，M=50-200×", "bold": True, "color": DARK, "size": Pt(16), "space": 6},
    {"text": "高反向偏压(100-300V) → 载流子碰撞电离链式倍增 → 信噪比改善M²倍。代表: Hamamatsu S8890, First Sensor", "level": 1, "size": Pt(15), "space": 8},
    {"text": "Level 3: SPAD单光子雪崩二极管 ---- Geiger模式，无限增益 ★", "bold": True, "color": ORANGE, "size": Pt(16), "space": 6},
    {"text": "偏压>击穿电压V_BR → 单光子即可触发自持雪崩 → 需要淬灭电路强制停止。PDE=5-30%, DCR=100-1000cps/μm², 死区时间10-100ns", "level": 1, "size": Pt(15), "space": 8},
    {"text": "Level 4: SiPM硅光电倍增管 ---- N个SPAD并联，光子数分辨 ★ 车载LiDAR核心", "bold": True, "color": GREEN, "size": Pt(16), "space": 6},
    {"text": "10³-10⁴个SPAD微单元独立工作→输出电流∝触发微元数→模拟光子计数。动态范围: 单光子→10⁴-10⁵photons/pulse, 响应时间<100ps", "level": 1, "size": Pt(15), "space": 8},
    {"text": "关键工艺: 3D堆叠(SPAD层+CMOS逻辑层TSV互连) → 消除键合寄生 → 提升填充因子和时间分辨率", "bold": True, "color": GREEN, "size": Pt(15), "space": 6},
    {"text": "代表产品: Sony IMX459(车规SPAD深度传感器) - Onsemi/SensL - Hamamatsu - 禾赛自研SiPM", "level": 1, "size": Pt(15)},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.0), items_all)
add_chart(s, "detector_chain.png", Inches(0.5), Inches(6.2), Inches(12.0), height=Inches(1.0))
add_page_num(s, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: LiDAR Measurement Circuit (1) -- TIA
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第四章  激光雷达测量电路（一）---- TIA跨阻放大器：从光电流到电压")

items_L = [
    {"text": "TIA (Transimpedance Amplifier) 核心功能", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "将SPAD/APD输出的μA-nA级微弱电流脉冲 → 可处理的电压信号", "level": 1, "size": Pt(16), "space": 8},
    {"text": "基本方程: V_out = -I_in × R_f", "bold": True, "level": 1, "size": Pt(15), "color": ORANGE, "space": 8},
    {"text": "典型参数: R_f=10-100kΩ, 1μA输入→10-100mV输出", "level": 1, "size": Pt(15), "space": 10},
    {"text": "", "size": Pt(4)},
    {"text": "高频稳定性设计 ★ 关键挑战", "bold": True, "color": MED, "size": Pt(19), "space": 10},
    {"text": "探测器结电容C_d(1-5pF) + R_f → RC极点 → 限制带宽", "level": 1, "size": Pt(15), "space": 6},
    {"text": "跨阻带宽要求: BW ≥ 0.35/t_r (1ns脉冲→BW≥350MHz!)", "level": 1, "size": Pt(15), "color": ORANGE, "space": 6},
    {"text": "并联补偿电容C_f → 闭环传递函数: H(s) = -R_f/(1+sR_fC_f)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "构成一阶低通滤波, -3dB带宽 f_c = 1/(2πR_fC_f)", "level": 2, "size": Pt(14), "space": 6},
    {"text": "设计要点: C_f需抵消C_d引入的极点并保证相位裕度>45°", "level": 1, "size": Pt(15), "color": GREEN, "space": 10},
    {"text": "先进TIA参数: 180nm CMOS工艺, BW=900MHz, 功耗1.5mW/ch", "bold": True, "level": 1, "size": Pt(15), "color": GREEN},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), items_L)

# Right: TIA formulas + circuit
add_formula(s, "TIA", Inches(6.8), Inches(1.5), width=Inches(5.5), height=Inches(0.5))
add_formula(s, "TIA_noise", Inches(6.8), Inches(2.1), width=Inches(5.5), height=Inches(0.5))
add_formula(s, "TIA_H", Inches(6.8), Inches(2.7), width=Inches(5.5), height=Inches(0.5))

items_R = [
    {"text": "TIA等效输入噪声三源模型:", "bold": True, "color": DARK, "size": Pt(15), "space": 8},
    {"text": "i²_n,total = i²_n,Rf + i²_n,op + v²_n,op×(1/Rf² + ω²C²_d)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "① 反馈电阻热噪声: 4kT/Rf (Johnson噪声)", "level": 1, "size": Pt(14), "space": 4},
    {"text": "② 运放输入电流噪声: i_n,op", "level": 1, "size": Pt(14), "space": 4},
    {"text": "③ 运放输入电压噪声经输入电容放大: v_n,op×ωC_d", "level": 1, "size": Pt(14), "space": 6},
    {"text": "信号增益∝Rf, 噪声中Rf项∝1/Rf → Rf↑→SNR↑", "level": 1, "size": Pt(15), "color": GREEN, "space": 6},
    {"text": "但Rf过大 → 带宽↓ + 振铃风险↑ → 设计折中", "level": 1, "size": Pt(15), "color": ORANGE},
]
add_textbox(s, Inches(6.5), Inches(3.4), Inches(6.3), Inches(3.2), items_R)

add_circuit(s, "tia_circuit", Inches(6.8), Inches(5.7), width=Inches(5.8), height=Inches(1.5))
add_page_num(s, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: LiDAR Measurement Circuit (2) -- TDC
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第四章  激光雷达测量电路（二）---- TDC时间数字转换器：皮秒级飞行时间测量")

items_L = [
    {"text": "TDC (Time-to-Digital Converter) 核心挑战", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "1cm测距精度 → Δt = 2×1cm/c ≈ 67ps", "bold": True, "level": 1, "size": Pt(15), "color": ORANGE, "space": 8},
    {"text": "单一时钟方案: 1GHz时钟 → 1ns分辨率 → 15cm误差 (不可接受!)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "TDC突破: 用电路延迟而非时钟周期来刻度时间", "level": 1, "size": Pt(15), "color": GREEN, "space": 12},
    {"text": "延迟线型TDC ---- 最成熟的车载方案", "bold": True, "color": MED, "size": Pt(19), "space": 10},
    {"text": "Start信号 → 进入Buffer链(τ_LSB≈10-20ps)逐级传播", "level": 1, "size": Pt(16), "space": 6},
    {"text": "Stop信号 → 触发D触发器阵列锁存各级状态", "level": 1, "size": Pt(16), "space": 6},
    {"text": "读出温度计码: 111...11000...0 → 细时间 Δt_fine = N×τ_LSB", "level": 1, "size": Pt(16), "space": 6},
    {"text": "粗计数器(Coarse) + 细时间(Fine) → 完整ToF", "level": 1, "size": Pt(16), "color": GREEN, "space": 6},
    {"text": "TI TDC7200: 55ps分辨率, 双通道, SPI接口, 车规级", "level": 2, "size": Pt(14), "space": 10},
    {"text": "Vernier游标型TDC ---- 更高精度(~2ps)", "bold": True, "color": MED, "size": Pt(19), "space": 8},
    {"text": "双延迟链: τ₁ vs τ₂ (微差1-2ps) → 游标原理放大 → 等效分辨率=|τ₁-τ₂|", "level": 1, "size": Pt(16)},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5), items_L)

# Right: TDC formulas + circuit
add_formula(s, "TDC", Inches(7.3), Inches(1.5), width=Inches(5.2), height=Inches(0.45))
add_formula(s, "TDC_jitter", Inches(7.3), Inches(2.1), width=Inches(5.2), height=Inches(0.5))

items_R = [
    {"text": "TDC关键性能指标:", "bold": True, "color": DARK, "size": Pt(15), "space": 8},
    {"text": "分辨率(LSB): 延迟链单级延迟=10-55ps", "level": 1, "size": Pt(15), "space": 6},
    {"text": "精度/抖动: σ_total = √(σ²_quant + σ²_DNL + σ²_INL + σ²_clk)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "量化误差: σ_quant = τ_LSB/√12 (均匀分布)", "level": 1, "size": Pt(14), "space": 6},
    {"text": "微分非线性DNL: 延迟单元失配 → ±0.5 LSB以内", "level": 1, "size": Pt(14), "space": 6},
    {"text": "积分非线性INL: 延迟线累积误差 → 需要校准", "level": 1, "size": Pt(14), "space": 6},
    {"text": "动态范围: 通常12-16bit → 10ns-10μs量程", "level": 1, "size": Pt(15), "space": 6},
    {"text": "死区时间: 转换+读出时间 → 影响最大脉冲率", "level": 1, "size": Pt(15), "space": 8},
    {"text": "先进集成方案:", "bold": True, "color": GREEN, "size": Pt(15), "space": 8},
    {"text": "SPAD阵列 + TDC阵列 3D堆叠SIP", "level": 1, "size": Pt(15), "space": 6},
    {"text": "每个SPAD微单元独立TDC → 并行ToF直方图", "level": 1, "size": Pt(15), "space": 6},
    {"text": "TCSPC时间相关单光子计数模式", "level": 1, "size": Pt(15), "color": ORANGE},
]
add_textbox(s, Inches(7.0), Inches(2.7), Inches(6.0), Inches(3.8), items_R)

add_circuit(s, "tdc_circuit", Inches(7.0), Inches(5.7), width=Inches(5.8), height=Inches(1.5))
add_page_num(s, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: LiDAR Measurement Circuit (3) -- SPAD Quenching
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第四章  激光雷达测量电路（三）---- SPAD淬灭与复位：让单光子探测器持续工作")

items_L = [
    {"text": "SPAD淬灭的必要性", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "Geiger模式下，单光子触发自持雪崩 → 电流持续增长 → 烧毁器件！", "level": 1, "size": Pt(16), "color": RED, "space": 8},
    {"text": "淬灭电路: 检测雪崩→迅速拉低偏压至V_BR以下→雪崩停止", "level": 1, "size": Pt(16), "space": 6},
    {"text": "复位电路: 雪崩停止后→恢复偏压至V_APD→准备下一次探测", "level": 1, "size": Pt(16), "space": 12},
    {"text": "无源淬灭 (Passive Quenching)", "bold": True, "color": DARK, "size": Pt(19), "space": 10},
    {"text": "串联大电阻R_Q(100kΩ-1MΩ) → 雪崩电流→R_Q压降→自动降低SPAD偏压", "level": 1, "size": Pt(15), "space": 6},
    {"text": "优点: 结构简单-面积小-零功耗", "level": 1, "size": Pt(15), "color": GREEN, "space": 6},
    {"text": "缺点: 死区时间50-500ns → 最大计数率<10MHz → 不满足车载需求", "level": 1, "size": Pt(15), "color": RED, "space": 6},
    {"text": "恢复时间: τ_recovery = R_Q×(C_SPAD+C_parasitic)", "level": 2, "size": Pt(14)},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.5), items_L)

items_R = [
    {"text": "主动淬灭 (Active Quenching) ★ 车载必选", "bold": True, "color": ORANGE, "size": Pt(19), "space": 12},
    {"text": "雪崩检测电路+快速MOS开关: 检测雪崩→主动拉低阳极→强制淬灭", "level": 1, "size": Pt(15), "space": 6},
    {"text": "优点: 死区时间2-10ns → 计数率>100MHz", "bold": True, "level": 1, "size": Pt(16), "color": GREEN, "space": 6},
    {"text": "缺点: 电路复杂-每SPAD需独立淬灭MOS-功耗增加", "level": 1, "size": Pt(15), "color": ORANGE, "space": 6},
    {"text": "先进方案: 混合淬灭(被动检测+主动关断+快速复位)", "level": 1, "size": Pt(15), "space": 10},
    {"text": "SiPM淬灭集成方案", "bold": True, "color": MED, "size": Pt(19), "space": 10},
    {"text": "每个SPAD微单元集成淬灭电阻+快速复位电路", "level": 1, "size": Pt(15), "space": 6},
    {"text": "3D堆叠: SPAD层(顶层) + 淬灭电路层(中间) + 数字处理层(底层)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "淬灭电阻典型值: 200kΩ-500kΩ (片上多晶硅电阻)", "level": 1, "size": Pt(14), "space": 6},
    {"text": "暗计数抑制: 相邻SPAD符合检测 → 有效消除热噪声误触发", "level": 1, "size": Pt(15), "color": GREEN, "space": 10},
    {"text": "代表产品: 禾赛自研SiPM+淬灭ASIC - Onsemi RB系列 - Hamamatsu MPPC", "bold": True, "level": 1, "size": Pt(15), "color": DGRAY},
]
add_textbox(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(5.0), items_R)

add_circuit(s, "spad_quenching_circuit", Inches(1.5), Inches(5.7), width=Inches(10.5), height=Inches(1.6))
add_page_num(s, 10)

# ═══════════════════════════════════════════════════════════
# SLIDE 11: LiDAR Applications
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第五章  激光雷达应用 ---- 自动驾驶感知、定位与建图")

items_L = [
    {"text": "自动驾驶感知任务 (L3/L4级)", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "① 3D目标检测: 车辆-行人-骑行者-锥桶-护栏 → 3D BBox+速度+航向角", "level": 1, "size": Pt(16), "space": 6},
    {"text": "② 可行驶区域分割: 地面/非地面分类 → 路径规划输入", "level": 1, "size": Pt(16), "space": 6},
    {"text": "③ 高精定位: 点云与HD Map匹配(ICP/NDT算法) → cm级定位精度", "level": 1, "size": Pt(16), "space": 6},
    {"text": "④ SLAM实时建图: LOAM/LIO-SAM框架 → LiDAR+IMU紧耦合", "level": 1, "size": Pt(16), "space": 6},
    {"text": "⑤ 目标追踪(MOT): 卡尔曼滤波+匈牙利算法 → 时序ID关联", "level": 1, "size": Pt(16), "space": 12},
    {"text": "典型传感器配置 (L4 Robotaxi)", "bold": True, "color": MED, "size": Pt(19), "space": 10},
    {"text": "1×主LiDAR(顶部, 120°H×25°V, 150m+)", "level": 1, "size": Pt(15), "space": 4},
    {"text": "2-4×补盲LiDAR(侧向, 短距) + 5-8×摄像头 + 5-8×毫米波雷达", "level": 1, "size": Pt(15), "space": 4},
    {"text": "1×组合导航(IMU+GNSS RTK) → 全车传感器冗余", "level": 1, "size": Pt(15), "color": GREEN},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5), items_L)

items_R = [
    {"text": "量产车型LiDAR搭载案例", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "蔚来ET7/ES7: 1×Innovusion Falcon(1550nm, 500m)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "理想L9/L8 Max: 1×禾赛AT128(128线, 905nm, 200m@10%)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "小鹏G9/P7i: 2×速腾聚创M1(MEMS固态, 150m)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "路特斯Eletre: 4×LiDAR(可升降式, 360°覆盖)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "华为问界M5/M7: 1×华为自研96线LiDAR", "level": 1, "size": Pt(15), "space": 10},
    {"text": "前装量产价格趋势", "bold": True, "color": ORANGE, "size": Pt(19), "space": 10},
    {"text": "2020: $1000-5000/颗 (机械式)", "level": 1, "size": Pt(16), "space": 4},
    {"text": "2023: $500-1000/颗 (MEMS/转镜)", "level": 1, "size": Pt(16), "space": 4},
    {"text": "2025E: $300-500/颗 (固态化+规模量产)", "level": 1, "size": Pt(16), "color": GREEN, "space": 4},
    {"text": "2030E: <$100/颗 (LiDAR-on-a-Chip)", "level": 1, "size": Pt(16), "color": ORANGE, "space": 10},
    {"text": "2025年中国前装LiDAR出货量: ~150万颗 (YoY+200%)", "bold": True, "level": 1, "size": Pt(16), "color": MED, "space": 10},
    {"text": "禾赛/速腾聚创/华为/图达通 合计>90%市场份额", "level": 1, "size": Pt(15)},
]
add_textbox(s, Inches(7.2), Inches(1.5), Inches(5.8), Inches(5.5), items_R)
add_page_num(s, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12: LiDAR Complete Signal Chain
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "激光雷达完整信号链 ---- 从光子到点云的全链路")

items_all = [
    {"text": "LiDAR信号处理全链路", "bold": True, "color": MED, "size": Pt(19), "space": 12},
    {"text": "", "size": Pt(4)},
    {"text": "【发射端】MCU/FPGA时序控制 → VCSEL驱动器(窄脉宽电流脉冲) → VCSEL阵列(905nm激光) → 发射光学(准直+扩束) → 自由空间传输 → 目标反射", "bold": True, "color": DARK, "size": Pt(16), "space": 10},
    {"text": "【接收端】接收光学(透镜汇聚) → 窄带滤光片(抑制背景光) → SPAD/SiPM(光电转换) → TIA(电流→电压) → VGA(可变增益放大) → 比较器/ADC(数字化) → TDC(皮秒计时)", "bold": True, "color": DARK, "size": Pt(16), "space": 10},
    {"text": "【数字处理】ToF直方图生成(TCSPC模式) → 峰值检测 → 距离计算R=½cΔt → 点云坐标转换(球→笛卡尔) → 多帧累积去噪 → 3D点云输出", "bold": True, "color": DARK, "size": Pt(16), "space": 10},
    {"text": "", "size": Pt(6)},
    {"text": "四大核心公式串联", "bold": True, "color": ORANGE, "size": Pt(19), "space": 10},
    {"text": "P_det = 1-exp(-∫η_PDE×Φ_ph dt) → I_SiPM = (N_fired/N_micro)×G×q → V_out = -I_in×R_f → Δt = N×τ_LSB → R = ½cΔt", "level": 1, "size": Pt(16), "color": ORANGE, "space": 8},
    {"text": "", "size": Pt(4)},
    {"text": "SPAD探测概率 → SiPM输出电流 → TIA跨阻放大 → TDC时间量化 → ToF距离计算", "level": 1, "size": Pt(15), "color": DGRAY, "space": 12},
    {"text": "关键系统指标: 探测概率Pd>99.7%@100m - 测距精度<2cm - 帧率>10Hz - 功耗<15W - 成本<$500(2025)", "bold": True, "color": GREEN, "size": Pt(16)},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.8), items_all)
add_page_num(s, 12)

# ═══════════════════════════════════════════════════════════
# SLIDE 13: Other Sensors (1) -- Camera & mmWave Radar
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "附录A  其他传感器技术概览（一）---- CMOS摄像机与毫米波雷达")

# LEFT: Camera
items_cam = [
    {"text": "CMOS摄像机 ---- 视觉感知核心", "bold": True, "color": MED, "size": Pt(16), "space": 10},
    {"text": "原理: 内光电效应 → 4T-APS像素 → CDS相关双采样消除kTC噪声", "level": 1, "size": Pt(14), "space": 5},
    {"text": "核心元器件: 微透镜阵列+彩色滤光片+PPD光电二极管+列级SAR-ADC", "level": 1, "size": Pt(14), "space": 5},
    {"text": "测量电路: CDS △V = V_signal - V_reset, 片上14bit ADC", "level": 1, "size": Pt(14), "space": 5},
    {"text": "车载挑战: HDR 120-140dB - LED闪烁抑制(LFM) - -40~+105°C", "level": 1, "size": Pt(14), "space": 5},
    {"text": "L2级5-8颗 → L3/L4级12-15+颗 → 市场$170亿(2024)", "level": 1, "size": Pt(14), "color": GREEN},
]
add_textbox(s, Inches(0.3), Inches(1.45), Inches(6.0), Inches(2.5), items_cam)
add_formula(s, "Iph", Inches(0.5), Inches(4.0), width=Inches(2.8))
add_formula(s, "CDS", Inches(3.5), Inches(4.0), width=Inches(2.8))
add_circuit(s, "cds_circuit", Inches(0.3), Inches(4.5), width=Inches(6.0), height=Inches(1.2))

# RIGHT: mmWave
items_mmw = [
    {"text": "毫米波雷达 ---- 全天候测距测速", "bold": True, "color": MED, "size": Pt(16), "space": 10},
    {"text": "原理: FMCW调频连续波, f_IF=(2B-R)/(c-T_c), f_d=2v_r/λ", "level": 1, "size": Pt(14), "space": 5},
    {"text": "核心元器件: 77GHz MMIC(PLL+VCO+PA+LNA+Mixer) - MIMO天线阵列", "level": 1, "size": Pt(14), "space": 5},
    {"text": "测量电路: 2D-FFT Range-Doppler处理 - 3D-FFT+Capon/MUSIC角度估计", "level": 1, "size": Pt(14), "space": 5},
    {"text": "4D成像雷达: 增加俯仰向天线 → 高度维度 → \"成像级\"点云", "level": 1, "size": Pt(14), "space": 5},
    {"text": "优势: 全天候-直接测速-200-400元/颗 - 国产化率<10%→~35%(2025)", "level": 1, "size": Pt(14), "color": GREEN},
]
add_textbox(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(2.5), items_mmw)
add_formula(s, "FMCW_TX", Inches(6.8), Inches(4.0), width=Inches(3.0))
add_formula(s, "FMCW_IF", Inches(10.0), Inches(4.0), width=Inches(3.0))
add_circuit(s, "fmcw_trx_circuit", Inches(6.8), Inches(4.5), width=Inches(6.0), height=Inches(1.2))

# Bottom comparison
items_comp = [
    {"text": "LiDAR vs Camera vs mmWave 互补关系:", "bold": True, "color": ORANGE, "size": Pt(16), "space": 6},
    {"text": "Camera(纹理/颜色/分类) + LiDAR(精确3D几何) + Radar(全天候/测速/direct速度) = 异构冗余 → L3+感知安全", "level": 1, "size": Pt(15), "color": DGRAY},
]
add_textbox(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.8), items_comp)
add_page_num(s, 13)

# ═══════════════════════════════════════════════════════════
# SLIDE 14: Other Sensors (2) -- Ultrasonic & IMU
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "附录B  其他传感器技术概览（二）---- 超声波雷达与MEMS IMU")

# LEFT: Ultrasonic
items_us = [
    {"text": "超声波雷达 ---- 短距泊车辅助", "bold": True, "color": MED, "size": Pt(16), "space": 10},
    {"text": "原理: PZT压电换能(40-58kHz), ToF: R=½v_sΔt, v_s=331.3+0.606T(°C)", "level": 1, "size": Pt(14), "space": 5},
    {"text": "核心元器件: PZT压电陶瓷-变压器-Elmos E524 ASIC(集成收发)", "level": 1, "size": Pt(14), "space": 5},
    {"text": "测量电路: 升压驱动(3.3V→100-300Vpp) → 接收链(LNA→BPF→VGA→比较器)", "level": 1, "size": Pt(14), "space": 5},
    {"text": "应用: UPA泊车-APA自动泊车-BSD盲区检测(0.3-5m)", "level": 1, "size": Pt(14), "space": 5},
    {"text": "局限: 探测距离<7m - 风噪/雨滴干扰 - 波束角60-120°", "level": 1, "size": Pt(14), "color": ORANGE},
]
add_textbox(s, Inches(0.3), Inches(1.45), Inches(6.0), Inches(2.5), items_us)
add_formula(s, "UltraR", Inches(0.5), Inches(4.0), width=Inches(2.8))
add_formula(s, "SoundV", Inches(3.5), Inches(4.0), width=Inches(2.8))
add_circuit(s, "ultrasonic_circuit", Inches(0.3), Inches(4.5), width=Inches(6.0), height=Inches(1.2))

# RIGHT: IMU
items_imu = [
    {"text": "MEMS IMU ---- 高频短时位姿粘合剂", "bold": True, "color": MED, "size": Pt(16), "space": 10},
    {"text": "加速度计: F=ma→位移→ΔC∝a, 开关电容读出+Σ-Δ ADC", "level": 1, "size": Pt(14), "space": 5},
    {"text": "陀螺仪: 科里奥利力F_c=-2m(Ω×v), 驱动模态→检测模态读出Ω", "level": 1, "size": Pt(14), "space": 5},
    {"text": "核心元器件: 硅质量块-梳齿电极-弹性悬臂梁-ASIC读出芯片", "level": 1, "size": Pt(14), "space": 5},
    {"text": "测量电路: 开关电容差分读出, ΔV≈V_ref×ΔC/(2C₀), Chopper抑制1/f噪声", "level": 1, "size": Pt(14), "space": 5},
    {"text": "误差模型: Allan方差辨识白噪声-零偏不稳定性-随机游走-速率斜坡", "level": 1, "size": Pt(14), "space": 5},
    {"text": "IMU角色: 高频(100-1000Hz)短时粘合剂, 从不单独使用→融合系统脊梁", "level": 1, "size": Pt(14), "color": GREEN},
]
add_textbox(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(2.5), items_imu)
add_formula(s, "Coriolis", Inches(6.8), Inches(4.0), width=Inches(3.0))
add_formula(s, "Allan", Inches(10.0), Inches(4.0), width=Inches(3.0))
add_circuit(s, "mems_readout_circuit", Inches(6.8), Inches(4.5), width=Inches(6.0), height=Inches(1.2))

add_page_num(s, 14)

# ═══════════════════════════════════════════════════════════
# SLIDE 15: Sensor Fusion
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第六章  多传感器融合 ---- BEV+Transformer：自动驾驶感知架构")

add_chart(s, "radar_chart.png", Inches(0.3), Inches(1.6), Inches(5.0), Inches(4.5))
add_chart(s, "sensor_comparison_table.png", Inches(0.2), Inches(5.8), Inches(13.0), height=Inches(1.2))

items_R = [
    {"text": "三级融合架构", "bold": True, "color": MED, "size": Pt(19), "space": 10},
    {"text": "① 数据级(Early): 原始数据层融合→RGB-D → 信息最丰富但计算量最大", "level": 1, "size": Pt(15), "space": 6},
    {"text": "② 特征级(Mid) ★ 产业主流: 各传感器独立编码→BEV空间统一融合", "level": 1, "size": Pt(15), "color": ORANGE, "space": 6},
    {"text": "③ 决策级(Late): 独立检测→目标列表→关联匹配 → 结构简单但信息损失", "level": 1, "size": Pt(15), "space": 12},
    {"text": "BEV+Transformer感知架构 (2023-25产业主流)", "bold": True, "color": GREEN, "size": Pt(19), "space": 10},
    {"text": "多视相机特征(LSS 2D→BEV) + LiDAR点云(体素化→PointPillars编码)", "level": 1, "size": Pt(15), "space": 6},
    {"text": "可变形注意力(Deformable Attention) → 多模态BEV特征自适应聚合", "level": 1, "size": Pt(15), "space": 6},
    {"text": "多任务头: 3D检测 + 车道线 + 可行驶区域 + 轨迹预测", "level": 1, "size": Pt(15), "space": 6},
    {"text": "端侧部署: TensorRT FP16/INT8量化 → Jetson Orin 25FPS@15W", "level": 1, "size": Pt(15), "color": GREEN, "space": 12},
    {"text": "传感器融合核心价值", "bold": True, "color": ORANGE, "size": Pt(19), "space": 10},
    {"text": "LiDAR提供精确3D几何 + Camera提供语义纹理 + Radar提供全天候速度", "level": 1, "size": Pt(16), "space": 4},
    {"text": "异构冗余 → 任一传感器失效不影响安全 → 功能安全ASIL-D要求", "level": 1, "size": Pt(16), "color": GREEN},
]
add_textbox(s, Inches(5.5), Inches(1.4), Inches(7.3), Inches(5.8), items_R)
add_page_num(s, 15)

# ═══════════════════════════════════════════════════════════
# SLIDE 16: Market Trends & Career
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "第七章  市场趋势与职业规划 ---- 激光雷达产业的黄金十年")

add_chart(s, "cost_curve.png", Inches(0.2), Inches(1.5), Inches(4.5), Inches(2.8))
add_chart(s, "market_players.png", Inches(4.9), Inches(1.5), Inches(4.2), Inches(2.8))
add_chart(s, "lidar_timeline.png", Inches(0.2), Inches(6.6), Inches(12.5), height=Inches(0.7))

# Trends
items_trend = [
    {"text": "五大技术趋势", "bold": True, "color": MED, "size": Pt(16), "space": 10},
    {"text": "① 固态化+芯片化: LiDAR-on-a-Chip(<$100, <10mm³) → 2030年量产", "level": 1, "size": Pt(14), "space": 4},
    {"text": "② FMCW LiDAR: 相干探测+直接测速 → 2026-27年样机 → L4 Robotaxi首发", "level": 1, "size": Pt(14), "space": 4},
    {"text": "③ 4D成像雷达: 192虚拟通道+俯仰测量 → $50-100 → L2+大规模前装", "level": 1, "size": Pt(14), "space": 4},
    {"text": "④ BEV+Transformer融合: 92%渗透率(2027E) → 端到端感知范式", "level": 1, "size": Pt(14), "space": 4},
    {"text": "⑤ 国产替代: 中国LiDAR占全球~40% → 核心芯片国产化率<15%→>80%(2030E)", "level": 1, "size": Pt(14), "color": GREEN},
]
add_textbox(s, Inches(0.5), Inches(4.5), Inches(8.5), Inches(2.5), items_trend)

# Career
items_career = [
    {"text": "五大核心岗位与薪资（一线城市）", "bold": True, "color": ORANGE, "size": Pt(16), "space": 10},
    {"text": "① 激光雷达硬件工程师(模拟IC/TIA/TDC): 应届25-40万, 3-5年50-80万", "level": 1, "size": Pt(14), "space": 3},
    {"text": "② 感知融合算法工程师(3D检测/BEV): 应届30-45万, 3-5年60-100万 ★最热门", "level": 1, "size": Pt(14), "color": ORANGE, "space": 3},
    {"text": "③ 自动驾驶嵌入式系统工程师(CUDA/TensorRT): 应届25-35万, 3-5年45-70万", "level": 1, "size": Pt(14), "space": 3},
    {"text": "④ 多传感器标定算法工程师(几何/优化): 应届22-35万, 3-5年40-65万", "level": 1, "size": Pt(14), "space": 3},
    {"text": "⑤ 传感器系统集成测试工程师(AEC-Q100): 应届18-28万, 3-5年35-55万", "level": 1, "size": Pt(14), "space": 3},
    {"text": "人才供需: 严重失衡 → 复合背景(电子+计算机+AI)最稀缺 → 黄金窗口期3-5年", "level": 0, "size": Pt(15), "color": GREEN, "bold": True},
]
add_textbox(s, Inches(9.2), Inches(1.5), Inches(3.8), Inches(5.5), items_career)
add_page_num(s, 16)

# ═══════════════════════════════════════════════════════════
# SLIDE 17: Key Formulas Summary
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "核心公式总览 ---- 激光雷达从物理到电路的数学框架")

formula_pairs = [
    (Inches(0.2), Inches(1.45), "LiDAR_ToF", Inches(5.5), "ToF测距基本方程: R = ½-c-Δt"),
    (Inches(0.2), Inches(2.05), "LidarEq", Inches(5.5), "激光雷达方程: 回波功率与SNR模型"),
    (Inches(0.2), Inches(2.65), "SPAD_prob", Inches(5.5), "SPAD探测概率: 非齐次泊松过程"),
    (Inches(0.2), Inches(3.25), "SiPM", Inches(5.5), "SiPM输出电流: 微单元并联光子数分辨"),
    (Inches(0.2), Inches(3.85), "TIA", Inches(5.5), "TIA跨阻放大: V_out = -I_in×R_f"),
    (Inches(0.2), Inches(4.45), "TIA_noise", Inches(5.5), "TIA三源噪声模型: 热噪声+电流+电压噪声"),
    (Inches(0.2), Inches(5.05), "TDC", Inches(5.5), "TDC延迟线计时: Δt = N_coarse×T_clk + N_fine×τ_LSB"),
    (Inches(0.2), Inches(5.65), "TDC_jitter", Inches(5.5), "TDC抖动分解: 量化+DNL+INL+时钟抖动的RSS"),
    (Inches(0.2), Inches(6.25), "DarkI", Inches(5.5), "暗电流温度模型: I_dark ∝ T^(3/2)×exp(-E_g/2kT)"),
]

for (left, top, name, w, label) in formula_pairs:
    add_formula(s, name, left, top, width=w, height=Inches(0.45))
    bx = s.shapes.add_textbox(Inches(6.0), top, Inches(7.0), Inches(0.36))
    tf = bx.text_frame; p = tf.paragraphs[0]
    p.text = label; p.font.size = Pt(13); p.font.color.rgb = DGRAY

add_page_num(s, 17)

# ═══════════════════════════════════════════════════════════
# SLIDE 18: Conclusion
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(s); add_title_bar(s, "总结与展望 ---- 激光雷达技术的核心结论")

items_all = [
    {"text": "本调研五大核心结论", "bold": True, "color": DARK, "size": Pt(20), "space": 10},
    {"text": "", "size": Pt(4)},
    {"text": "1. 激光雷达是L3+自动驾驶不可或缺的3D几何感知核心传感器", "bold": True, "size": Pt(19), "color": MED, "space": 6},
    {"text": "ToF(主流)与FMCW(下一代)两大测距路线并存，VCSEL+SPAD/SiPM固态化方案推动成本从$5000→$100。", "level": 1, "size": Pt(15), "space": 8},
    {"text": "2. SPAD/SiPM单光子探测器是实现远距离探测的关键器件", "bold": True, "size": Pt(19), "color": MED, "space": 6},
    {"text": "PIN→APD→SPAD→SiPM 灵敏度逐级百万倍提升，3D堆叠+TSV工艺是降本提效的核心路径。", "level": 1, "size": Pt(15), "space": 8},
    {"text": "3. 精密测量电路(TIA+TDC+淬灭)是信号链的决定性环节", "bold": True, "size": Pt(19), "color": MED, "space": 6},
    {"text": "TIA将μA电流转为电压，TDC以皮秒分辨率量化飞行时间，主动淬灭电路保障SPAD持续工作----三者缺一不可。", "level": 1, "size": Pt(15), "space": 8},
    {"text": "4. 多传感器融合(BEV+Transformer)是L3+感知的标准架构", "bold": True, "size": Pt(19), "color": MED, "space": 6},
    {"text": "LiDAR(几何)+Camera(语义)+Radar(全天候速度)+IMU(高频惯性)异构互补→功能安全冗余。", "level": 1, "size": Pt(15), "space": 8},
    {"text": "5. 2025-2030是激光雷达产业的黄金十年", "bold": True, "size": Pt(19), "color": MED, "space": 6},
    {"text": "中国前装出货150万颗(2025)→CAGR+200%，LiDAR-on-a-Chip <$100(2030E)，复合型人才严重供不应求。", "level": 1, "size": Pt(15), "space": 10},
    {"text": "核心能力公式: 物理电子学基础 + 嵌入式系统实践 + 深度学习算法能力 = 综合竞争力", "bold": True, "size": Pt(19), "color": ORANGE},
]
add_textbox(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5), items_all)
add_page_num(s, 18)

# ═══════════════════════════════════════════════════════════
# SLIDE 19: Thanks
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(s)

for y, h, c in [(0, 0.08, ORANGE), (7.42, 0.08, MED)]:
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), SW, Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

box = s.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.3), Inches(1.8))
tf = box.text_frame
p = tf.paragraphs[0]; p.text = "感谢观看！欢迎交流指正"
p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "激光雷达技术深度解析 ---- 从物理原理到自动驾驶应用"
p2.font.size = Pt(16); p2.font.color.rgb = ACCLR; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
p3 = tf.add_paragraph()
p3.text = "广东工业大学 - 传感器技术与应用 - 先进传感器调研"
p3.font.size = Pt(14); p3.font.color.rgb = ACCLR; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(6)

refs = [
    "主要参考文献 (共20+篇, 覆盖原理/器件/电路/算法/市场各方向)",
    "",
    "■ 方吉鑫等. 硅光电倍增管前端放大电路架构研究进展[J]. 激光技术, 2026.",
    "■ 林远启等. 集成化多线列激光雷达模拟前端微组件设计[J]. 光电工程, 2021.",
    "■ 于航. 基于激光雷达的3D目标检测研究综述[J]. 汽车文摘, 2024(2):18-27.",
    "■ 李卓达, 查云飞. 智能汽车多源融合环境感知系统概述[J]. 汽车文摘, 2025(4):1-11.",
    "■ Poulton C V et al. Optical Beamforming for Solid-State Automotive LIDAR[R]. UC Berkeley, 2023.",
    "■ Li Y et al. Lidar for Autonomous Driving[J]. IEEE Signal Processing Magazine, 2020.",
    "■ Lang A H et al. PointPillars: Fast Encoders for Object Detection from Point Clouds[C]. CVPR, 2019.",
    "■ A Review of Multi-Sensor Fusion in Autonomous Driving[J]. Sensors, 2025, 25(19):6033.",
    "■ Wang X Y et al. A Review of Image and Point Cloud Fusion in AD[C]. LNICST, 2024, 554:62-73.",
    "■ Yole Intelligence. LiDAR for Automotive 2024[R]; 高工智能汽车研究院, 2025.",
    "■ 禾赛科技-速腾聚创-华为-图达通-Sony-Onsemi-Hamamatsu 等企业产品技术资料.",
]
ref_items = []
for r in refs:
    if r and r[0] == "■":
        ref_items.append({"text": r, "size": Pt(14), "color": WHITE, "bold": False})
    elif r:
        ref_items.append({"text": r, "size": Pt(14), "color": WHITE, "bold": True})
    else:
        ref_items.append({"text": "", "size": Pt(8)})
add_textbox(s, Inches(1), Inches(3.0), Inches(11.3), Inches(3.8), ref_items)

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
outpath = os.path.join(BASE, "自动驾驶汽车先进传感器技术综述_LiDAR专题.pptx")
prs.save(outpath)
print(f"PPT saved: {outpath}")
print(f"Slides: {len(prs.slides)}")
